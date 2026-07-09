"""`compare-office` skill — xlsx / pptx / docx from one ComparisonResult (P7).

Pure functions, no fetch/recompute/network. Every deliverable is driven from the
SAME payload so all four (with compare-viz) tell one story. Office libraries are
imported lazily so the package imports without the `dev` extra. Traceability:
each asset's `source_url` is carried into the workbook and the narrative.
"""

from __future__ import annotations

import io

from ...models import ComparisonResult


def _closes_by_ticker(result: ComparisonResult) -> dict[str, dict[str, float]]:
    return {s.ticker: {b.date: b.close for b in s.bars} for s in result.series}


def to_xlsx(result: ComparisonResult) -> bytes:
    """Backtest 底稿: Overview / Prices / Returns / Correlation / Backtests / Sources."""
    from openpyxl import Workbook  # noqa: PLC0415

    wb = Workbook()
    dates = result.aligned_dates
    cbt = _closes_by_ticker(result)
    tickers = result.tickers

    ov = wb.active
    ov.title = "Overview"
    ov.append([result.title, f"{result.range[0]}..{result.range[1]}",
               f"{len(dates)} common trading days"])
    ov.append([])
    ov.append(["ticker", "start", "end", "total_return", "cagr", "annual_vol",
               "sharpe", "max_drawdown", "dd_start", "dd_end"])
    for m in result.metrics:
        ov.append([m.ticker, m.start_price, m.end_price, m.total_return, m.cagr,
                   m.annual_vol, m.sharpe, m.max_drawdown,
                   m.drawdown_window[0], m.drawdown_window[1]])

    pr = wb.create_sheet("Prices")
    pr.append(["date", *tickers])
    for d in dates:
        pr.append([d, *[cbt[t][d] for t in tickers]])

    rt = wb.create_sheet("Returns")
    rt.append(["date", *[f"{t}_ret" for t in tickers]])
    for i in range(1, len(dates)):
        row: list[object] = [dates[i]]
        for t in tickers:
            prev, cur = cbt[t][dates[i - 1]], cbt[t][dates[i]]
            row.append(cur / prev - 1.0 if prev else 0.0)
        rt.append(row)

    co = wb.create_sheet("Correlation")
    co.append(["pair", "pearson", "rolling_window"])
    for c in result.correlations:
        co.append([f"{c.ticker_a} vs {c.ticker_b}", c.pearson, c.rolling_window])
    co.append([])
    if result.correlations:
        c0 = result.correlations[0]
        co.append(["date", f"rolling_corr[{c0.ticker_a},{c0.ticker_b}]"])
        for d, v in c0.rolling:
            co.append([d, v])

    bt = wb.create_sheet("Backtests")
    bt.append(["name", "weights", "rebalance", "cost_bps", "total_return", "cagr",
               "annual_vol", "sharpe", "max_drawdown", "n_rebalances", "cost_drag"])
    for b in result.backtests:
        weights = ", ".join(f"{t}:{w:.2f}" for t, w in b.config.weights.items())
        bt.append([b.config.name, weights, b.config.rebalance, b.config.cost_bps,
                   b.total_return, b.cagr, b.annual_vol, b.sharpe, b.max_drawdown,
                   b.n_rebalances, b.cost_drag])
    bt.append([])
    if result.backtests:
        bt.append(["date", *[b.config.name for b in result.backtests]])
        eq_maps = [dict(b.equity_curve) for b in result.backtests]
        for d in dates:
            bt.append([d, *[m.get(d, "") for m in eq_maps]])

    sr = wb.create_sheet("Sources")
    sr.append(["ticker", "source", "source_url", "fetched_at"])
    for s in result.series:
        sr.append([s.ticker, s.prov.source, s.prov.source_url, s.prov.fetched_at])

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def to_pptx(result: ComparisonResult) -> bytes:
    """Comparison deck: context → per-asset performance → backtests → sources."""
    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Inches, Pt  # noqa: PLC0415

    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[5]

    def _title_slide(text: str) -> None:
        prs.slides.add_slide(title_layout).shapes.title.text = text

    def _bullets(title: str, lines: list[str]) -> None:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(6)).text_frame
        tb.word_wrap = True
        tb.text = title
        tb.paragraphs[0].font.size = Pt(26)
        for ln in lines:
            p = tb.add_paragraph()
            p.text = ln
            p.font.size = Pt(14)

    _title_slide(f"{result.title} — 多资产对比与策略回测")
    _bullets("Context", [
        f"Assets: {', '.join(result.tickers)}",
        f"Range: {result.range[0]} .. {result.range[1]}",
        f"Common trading days: {len(result.aligned_dates)}",
        "Deterministic, LLM-free; 252-day annualization, risk-free = 0.",
    ])
    _bullets("Per-asset performance (buy & hold)", [
        f"{m.ticker}: total {m.total_return:+.1%} · CAGR {m.cagr:+.1%} · "
        f"vol {m.annual_vol:.1%} · Sharpe {m.sharpe:.2f} · max DD {m.max_drawdown:.1%}"
        for m in result.metrics
    ] or ["(none)"])
    _bullets("Strategy backtests", [
        f"{b.config.name} [{b.config.rebalance}]: total {b.total_return:+.1%} · "
        f"CAGR {b.cagr:+.1%} · max DD {b.max_drawdown:.1%} · "
        f"{b.n_rebalances} rebalances · cost drag {b.cost_drag:.2%}"
        for b in result.backtests
    ] or ["(none)"])
    _bullets("Correlation", [
        f"{c.ticker_a} vs {c.ticker_b}: pearson {c.pearson:+.2f} "
        f"(daily returns, {c.rolling_window}d rolling series in xlsx)"
        for c in result.correlations
    ] or ["(none)"])
    _bullets("Sources", [
        f"{s.ticker} ({s.prov.source}): {s.prov.source_url}" for s in result.series
    ] or ["(none)"])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def to_docx(result: ComparisonResult) -> bytes:
    """Narrative comparison report over the same records, source-cited."""
    from docx import Document  # noqa: PLC0415

    doc = Document()
    doc.add_heading(f"{result.title} — Multi-Asset Comparison & Backtest", level=0)
    doc.add_paragraph(
        f"Assets {', '.join(result.tickers)} over {result.range[0]} .. {result.range[1]}. "
        f"Generated {result.generated_at}. {len(result.aligned_dates)} common trading days."
    )

    doc.add_heading("Methodology", level=1)
    doc.add_paragraph(
        "All statistics are deterministic and reproduce from the same inputs. Assets are "
        "compared on the intersection of their trading dates, so a 7-day-a-week asset and a "
        "weekday-only asset are measured on the days both traded. Returns are annualized on a "
        "252-day basis with a zero risk-free rate. Backtests hold target weights and, on each "
        "rebalance date, reset to those weights while charging the traded turnover a "
        "transaction cost — so rebalancing's frictional drag is modelled, not assumed free."
    )

    doc.add_heading("Per-asset performance", level=1)
    for m in result.metrics:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(f"{m.ticker}: ").bold = True
        p.add_run(
            f"total {m.total_return:+.1%}, CAGR {m.cagr:+.1%}, vol {m.annual_vol:.1%}, "
            f"Sharpe {m.sharpe:.2f}, max drawdown {m.max_drawdown:.1%} "
            f"({m.drawdown_window[0]} → {m.drawdown_window[1]})."
        )

    doc.add_heading("Strategy backtests", level=1)
    if result.backtests:
        for b in result.backtests:
            doc.add_heading(b.config.name, level=2)
            weights = ", ".join(f"{t} {w:.0%}" for t, w in b.config.weights.items())
            doc.add_paragraph(
                f"Weights {weights}; rebalance {b.config.rebalance}; "
                f"cost {b.config.cost_bps:.0f} bps. "
                f"Total return {b.total_return:+.1%}, CAGR {b.cagr:+.1%}, vol {b.annual_vol:.1%}, "
                f"Sharpe {b.sharpe:.2f}, max drawdown {b.max_drawdown:.1%}. "
                f"{b.n_rebalances} rebalances cost {b.cost_drag:.2%} of capital in frictions."
            )
    else:
        doc.add_paragraph("No backtest strategies were configured for this comparison.")

    doc.add_heading("Correlation", level=1)
    if result.correlations:
        for c in result.correlations:
            doc.add_paragraph(
                f"{c.ticker_a} vs {c.ticker_b}: daily-return Pearson correlation "
                f"{c.pearson:+.2f} over the window (a {c.rolling_window}-day rolling series "
                "is in the workbook).",
                style="List Bullet",
            )
    else:
        doc.add_paragraph("Only one asset was supplied; no pairwise correlation was computed.")

    doc.add_heading("Sources", level=1)
    for s in result.series:
        doc.add_paragraph(f"{s.ticker} ({s.prov.source}): {s.prov.source_url}", style="List Bullet")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()
