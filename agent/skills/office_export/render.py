"""`office-export` skill — xlsx / pptx / docx from one AnalysisResult (docs/phases/P4).

Pure functions, no fetch/recompute/network. Every deliverable is driven from the
SAME payload so all four tell the same story. Office libraries are imported
lazily inside each function so the package imports without the `full`/`dev`
extra; only calling an exporter needs the dependency.

Traceability: every xlsx row carries a source_url (P-INV-1 propagated to the sheet).
"""

from __future__ import annotations

import io

from ...models import AnalysisResult

_BAR_HEADERS = ["date", "open", "high", "low", "close", "volume", "source", "source_url"]
_INF_HEADERS = ["date", "kind", "significance", "price", "detector", "evidence_bars", "window"]
_EVT_HEADERS = ["date", "category", "impact", "title", "rationale", "news_refs", "source_url"]
_ALN_HEADERS = ["inflection_date", "kind", "lag_days", "confidence", "event_titles", "explanation"]


def to_xlsx(result: AnalysisResult) -> bytes:
    """Backtest 底稿: Bars / Inflections / Events / Alignments tabs, each row traceable."""
    from openpyxl import Workbook  # noqa: PLC0415

    wb = Workbook()
    bars = wb.active
    bars.title = "Bars"
    bars.append(_BAR_HEADERS)
    for b in result.series.bars:
        bars.append([b.date, b.open, b.high, b.low, b.close, b.volume,
                     b.prov.source, b.prov.source_url])

    inf = wb.create_sheet("Inflections")
    inf.append(_INF_HEADERS)
    for i in result.inflections:
        inf.append([i.date, i.kind.value, i.significance, i.price, i.detector,
                    ", ".join(i.evidence_bars), f"{i.window[0]}..{i.window[1]}"])

    evt = wb.create_sheet("Events")
    evt.append(_EVT_HEADERS)
    for e in result.events:
        evt.append([e.date, e.category, e.impact.value, e.title, e.rationale,
                    ", ".join(e.news_refs), e.prov.source_url])

    aln = wb.create_sheet("Alignments")
    aln.append(_ALN_HEADERS)
    for a in result.alignments:
        aln.append([a.inflection.date, a.inflection.kind.value, a.lag_days, a.confidence,
                    " | ".join(ev.title for ev in a.events), a.explanation])

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def to_pptx(result: AnalysisResult) -> bytes:
    """Decision-framework deck: context → inflections → aligned events → sources."""
    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Inches, Pt  # noqa: PLC0415

    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[5]

    def _title_slide(text: str) -> None:
        s = prs.slides.add_slide(title_layout)
        s.shapes.title.text = text

    def _bullets(title: str, lines: list[str]) -> None:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(6)).text_frame
        tb.word_wrap = True
        tb.text = title
        tb.paragraphs[0].font.size = Pt(26)
        for ln in lines:
            p = tb.add_paragraph()
            p.text = ln
            p.font.size = Pt(14)

    _title_slide(f"{result.ticker} — 行情与 AI 事件对齐分析")
    _bullets(
        "Context",
        [
            f"Ticker: {result.ticker}",
            f"Range: {result.range[0]} .. {result.range[1]}",
            f"Bars analyzed: {len(result.series.bars)}",
            f"Inflections detected: {len(result.inflections)}",
            f"Curated events: {len(result.events)}  ·  Alignments: {len(result.alignments)}",
        ],
    )
    _bullets(
        "Detected inflections",
        [f"{i.date}  {i.kind.value}  (significance {i.significance:.2f})"
         for i in result.inflections[:12]] or ["(none)"],
    )
    if result.alignments:
        for a in result.alignments[:12]:
            titles = "; ".join(ev.title for ev in a.events)
            _bullets(
                f"{a.inflection.date} · {a.inflection.kind.value}",
                [
                    f"Aligned event(s): {titles}",
                    f"Impact: {', '.join(ev.impact.value for ev in a.events)}",
                    f"Lag: {a.lag_days}d  ·  Confidence: {a.confidence:.0%}",
                    f"Why: {a.explanation}",
                ],
            )
    else:
        _bullets("Aligned events", ["No industry events were sourced for this window."])
    _bullets(
        "Sources",
        [f"{e.title} — {', '.join(e.news_refs)}" for e in result.events] or ["(none)"],
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def to_docx(result: AnalysisResult) -> bytes:
    """Narrative strategy report over the same records, citing news_refs."""
    from docx import Document  # noqa: PLC0415

    doc = Document()
    doc.add_heading(f"{result.ticker} — Market & AI-Event Analysis", level=0)
    doc.add_paragraph(
        f"Range {result.range[0]} .. {result.range[1]}. Generated {result.generated_at}. "
        f"{len(result.series.bars)} daily bars analyzed."
    )

    doc.add_heading("Methodology", level=1)
    det = result.inflections[0].detector if result.inflections else "n/a"
    doc.add_paragraph(
        "Price inflections are detected by a deterministic change-point algorithm "
        f"({det}); the same inputs reproduce the same points. Industry events are curated "
        "and aligned to inflections within a bounded time window; every conclusion links "
        "back to a cited source."
    )

    doc.add_heading("Findings", level=1)
    if result.alignments:
        for a in result.alignments:
            doc.add_heading(f"{a.inflection.date} — {a.inflection.kind.value}", level=2)
            doc.add_paragraph(a.explanation)
            for ev in a.events:
                p = doc.add_paragraph(style="List Bullet")
                p.add_run(f"{ev.title} [{ev.impact.value}] ({ev.date}) — ").bold = True
                p.add_run(ev.rationale)
                doc.add_paragraph("Sources: " + ", ".join(ev.news_refs), style="List Bullet 2"
                                  if _has_style(doc, "List Bullet 2") else "List Bullet")
    else:
        doc.add_paragraph("No industry events were sourced for this window; "
                          "price inflections are reported without event annotations.")

    doc.add_heading("Sources", level=1)
    for e in result.events:
        doc.add_paragraph(f"{e.title}: {', '.join(e.news_refs)}", style="List Bullet")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _has_style(doc: object, name: str) -> bool:
    try:
        return name in [s.name for s in doc.styles]  # type: ignore[attr-defined]
    except Exception:  # noqa: BLE001
        return False
