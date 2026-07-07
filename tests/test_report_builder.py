"""P4 tests for the deliverable renderers (T4.1–T4.8)."""

from __future__ import annotations

import io
import re
import socket
from pathlib import Path

import pytest

from agent.models import (
    Alignment,
    AnalysisResult,
    Bar,
    CuratedEvent,
    Impact,
    Inflection,
    InflectionKind,
    PriceSeries,
    Provenance,
)
from agent.skills.kline_viz import render_html
from agent.skills.office_export import to_docx, to_pptx, to_xlsx
from agent.subagents.report_builder import report_builder

FROZEN = "2025-01-01T00:00:00Z"
_PROV = Provenance("yahoo", "https://finance.yahoo.com/quote/NVDA/history", FROZEN, "NVDA")
_URL = "https://openai.com/blog/chatgpt"

# External http(s) in <script src>/<link href> — the thing P-INV-5 forbids.
_EXTERNAL = re.compile(r"<(?:script|link)\b[^>]*\b(?:src|href)\s*=\s*['\"]https?://", re.I)
_SECRET = re.compile(r"sk-[A-Za-z0-9]{16,}|AKIA[0-9A-Z]{16}|AIza[0-9A-Za-z_\-]{20,}")


def _bars() -> list[Bar]:
    rows = [
        ("2023-05-22", 30.0, 31.0, 29.5, 30.5, 100),
        ("2023-05-23", 30.5, 31.2, 30.1, 30.9, 110),
        ("2023-05-24", 30.9, 31.5, 30.6, 31.0, 120),
        ("2023-05-25", 31.0, 38.5, 30.9, 37.98, 900),
        ("2023-05-26", 38.0, 39.0, 37.5, 38.8, 400),
    ]
    return [Bar(d, o, h, lo, c, v, _PROV) for d, o, h, lo, c, v in rows]


def _result(*, with_events: bool = True, malicious: bool = False) -> AnalysisResult:
    bars = _bars()
    series = PriceSeries("NVDA", bars, _PROV)
    infl = Inflection("2023-05-25", InflectionKind.BREAKOUT_UP, 1.0, 37.98,
                      ("2023-05-20", "2023-05-30"), ["2023-05-24", "2023-05-25"],
                      "pelt:rbf:pen=8")
    if not with_events:
        return AnalysisResult("NVDA", ("2023-05-01", "2023-05-31"), series, [infl], [], [], FROZEN)

    title = "<script>alert(1)</script>" if malicious else "NVIDIA blowout Q1 FY24 earnings"
    prov = Provenance("curated", _URL, FROZEN, "NVDA", note="curated")
    event = CuratedEvent(title, "2023-05-24", "earnings", Impact.HIGH,
                         "Record data-center revenue guidance drove the gap up.", [_URL], prov)
    align = Alignment(infl, [event], 1, 0.95, "Breakout follows the earnings beat.")
    return AnalysisResult(
        "NVDA", ("2023-05-01", "2023-05-31"), series, [infl], [event], [align], FROZEN
    )


# --- T4.1 self-contained HTML (P-INV-5) ----------------------------------- #
def test_html_is_self_contained() -> None:
    html = render_html(_result())
    assert not _EXTERNAL.search(html), "external script/link src found"
    assert "echarts" in html.lower()
    assert len(html) > 500_000, "ECharts bundle does not appear inlined"


# --- T4.2 traceability propagated to the UI (P-INV-1/P-INV-2) ------------- #
def test_event_markers_link_to_real_sources() -> None:
    result = _result()
    html = render_html(result)
    for ev in result.events:
        for url in ev.news_refs:
            assert url in html  # the clickable source URL is present in the drill-down


# --- T4.3 output escaping (stored-XSS prevention) ------------------------- #
def test_malicious_title_is_escaped() -> None:
    html = render_html(_result(malicious=True))
    assert "&lt;script&gt;alert(1)&lt;/script&gt;" in html   # escaped form present
    assert "<script>alert(1)</script>" not in html           # raw tag NOT present


# --- T4.4 xlsx structure + traceable rows --------------------------------- #
def test_xlsx_tabs_and_source_urls() -> None:
    from openpyxl import load_workbook  # noqa: PLC0415

    wb = load_workbook(io.BytesIO(to_xlsx(_result())))
    assert wb.sheetnames == ["Bars", "Inflections", "Events", "Alignments"]
    events = wb["Events"]
    header = [c.value for c in events[1]]
    url_col = header.index("source_url")
    rows = list(events.iter_rows(min_row=2, values_only=True))
    assert rows, "Events tab is empty"
    for r in rows:
        assert r[url_col] and str(r[url_col]).startswith("http")


# --- T4.5 pptx + docx open and carry the story ---------------------------- #
def test_pptx_opens_with_ticker_and_events() -> None:
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(io.BytesIO(to_pptx(_result())))
    text = "\n".join(
        shape.text_frame.text
        for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
    )
    assert "NVDA" in text
    assert "NVIDIA blowout Q1 FY24 earnings" in text


def test_docx_opens_with_ticker_and_events() -> None:
    from docx import Document  # noqa: PLC0415

    doc = Document(io.BytesIO(to_docx(_result())))
    text = "\n".join(p.text for p in doc.paragraphs)
    assert "NVDA" in "\n".join(h.text for h in doc.paragraphs) or "NVDA" in text
    assert "NVIDIA blowout Q1 FY24 earnings" in text
    assert _URL in text  # sources cited


# --- T4.6 no-secret over produced artifacts (P-INV-4) --------------------- #
def test_no_secret_in_artifacts(tmp_path: Path) -> None:
    paths = report_builder(_result(), ["html", "xlsx", "pptx", "docx"], out_dir=str(tmp_path))
    assert len(paths) == 4
    for p in paths:
        text = Path(p).read_text(encoding="utf-8", errors="ignore")
        assert not _SECRET.search(text), f"key-shaped string in {p}"


# --- T4.7 purity: deterministic + no network ------------------------------ #
def test_html_render_is_deterministic() -> None:
    assert render_html(_result()) == render_html(_result())


def test_exporters_touch_no_network(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    def _boom(*_a: object, **_k: object) -> None:
        raise AssertionError("exporter attempted network access")

    monkeypatch.setattr(socket, "socket", _boom)
    paths = report_builder(_result(), ["html", "xlsx", "pptx", "docx"], out_dir=str(tmp_path))
    assert len(paths) == 4
    assert all(Path(p).exists() for p in paths)


# --- T4.8 graceful degradation when no events ----------------------------- #
def test_empty_events_still_renders(tmp_path: Path) -> None:
    result = _result(with_events=False)
    html = render_html(result)
    assert "No industry events were sourced" in html
    assert "candlestick" in html  # K-line still built
    # Office files still open with empty event bodies.
    from openpyxl import load_workbook  # noqa: PLC0415

    wb = load_workbook(io.BytesIO(to_xlsx(result)))
    assert wb.sheetnames == ["Bars", "Inflections", "Events", "Alignments"]
    paths = report_builder(result, ["html", "xlsx", "pptx", "docx"], out_dir=str(tmp_path))
    assert len(paths) == 4


# --- report_builder no-op on unknown outputs ------------------------------ #
def test_unknown_outputs_are_ignored(tmp_path: Path) -> None:
    assert report_builder(_result(), ["nope", "csv"], out_dir=str(tmp_path)) == []
